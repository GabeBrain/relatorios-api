#!/usr/bin/env python3
"""
Gerador de ESTUDO VOCACIONAL SINTÉTICO para testar o Corretor v3 barato.

Produz um .pptx pequeno (~15 slides, < 5 MB) com ERROS INJETADOS conhecidos +
um GABARITO JSON, para medir recall/FP do pipeline sem gastar ~R$ 2 de um estudo real.

As restrições de detecção foram lidas do código (jul/2026) e estão codificadas aqui:
  - Títulos casam os regexes de seção de  src/features/corretor/lib/audit/pptx-to-ir.ts
  - Tabelas-imagem respeitam a heurística de  lib/v3/table-images.ts
    (PNG 15-500 KB, largura >= 700px, altura 200-1300px, aspecto >= 1.2, seção numérica)
  - Regras DET de  lib/audit/ir-rules.ts  (LEFTOVER_NOTE, RADII, ABSOLUTE_SUM nativa)

Camadas de verificação do gabarito:
  DET    -> checável de graça (triagem no navegador / teste vitest pipeline-det.test.ts)
  IA     -> exige passe de texto (analyze-text-batch) — verificar no app (~R$ 0,02)
  VISION -> exige passe de visão (analyze-table-image) — verificar no app (~R$ 0,20)

Uso:
  python gera_estudo_teste.py --cidade Itajai --out teste_corretor.pptx
  python gera_estudo_teste.py --sem-secao entorno   # variante p/ STRUCTURE_MISSING

Requer: python-pptx, Pillow.
"""
import argparse
import io
import json
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

# ---------------------------------------------------------------- fonte p/ imagens
def _load_font(size):
    for path in (
        r"C:\Windows\Fonts\arial.ttf",
        r"C:\Windows\Fonts\segoeui.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ):
        if os.path.exists(path):
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def render_table_png(header, rows, total_row, width=1400, row_h=64):
    """Desenha uma tabela como PNG (>= 700px largura, aspecto >= 1.2, 15-500 KB).
    header: list[str]; rows: list[list[str]]; total_row: list[str] | None."""
    ncols = len(header)
    all_rows = [header] + rows + ([total_row] if total_row else [])
    height = row_h * len(all_rows) + 24
    # garante aspecto >= 1.2 e altura na faixa 200-1300
    height = max(220, min(1300, height))
    img = Image.new("RGB", (width, height), "white")
    d = ImageDraw.Draw(img)
    fh = _load_font(24)
    fb = _load_font(22)

    col_w = width / ncols
    y = 12
    for ri, row in enumerate(all_rows):
        is_header = ri == 0
        is_total = total_row is not None and ri == len(all_rows) - 1
        if is_header:
            d.rectangle([0, y, width, y + row_h], fill=(31, 78, 121))
        elif is_total:
            d.rectangle([0, y, width, y + row_h], fill=(221, 235, 247))
        elif ri % 2 == 0:
            d.rectangle([0, y, width, y + row_h], fill=(244, 247, 251))
        for ci, cell in enumerate(row):
            x = ci * col_w + 14
            color = "white" if is_header else (0, 0, 0)
            font = fh if is_header else fb
            d.text((x, y + row_h / 2 - 13), str(cell), fill=color, font=font)
        # linhas de grade
        d.line([0, y, width, y], fill=(200, 200, 200))
        y += row_h
    d.line([0, y, width, y], fill=(200, 200, 200))
    for ci in range(ncols + 1):
        d.line([ci * col_w, 12, ci * col_w, y], fill=(200, 200, 200))

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


# ---------------------------------------------------------------- helpers de slide
BLUE = RGBColor(0x1F, 0x4E, 0x79)


def add_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # "Title Only"
    slide.shapes.title.text = title
    return slide


def add_text(slide, text, top=2.2, left=0.6, width=8.0, height=1.4, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size)
    return box


def add_legend(slide, linhas, top=2.2, left=0.6):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(4.5), Inches(2.5))
    tf = box.text_frame
    tf.text = linhas[0]
    for extra in linhas[1:]:
        p = tf.add_paragraph()
        p.text = extra
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(20)
    return box


def add_native_table(slide, header, rows, total_row=None, top=2.2, left=0.6, width=8.5):
    all_rows = [header] + rows + ([total_row] if total_row else [])
    n, m = len(all_rows), len(header)
    gf = slide.shapes.add_table(n, m, Inches(left), Inches(top), Inches(width), Inches(0.4 * n))
    tbl = gf.table
    for ri, row in enumerate(all_rows):
        for ci, val in enumerate(row):
            tbl.cell(ri, ci).text = str(val)
    return tbl


def add_table_image(slide, header, rows, total_row=None, top=1.8, left=0.6, width_in=8.4):
    buf = render_table_png(header, rows, total_row)
    slide.shapes.add_picture(buf, Inches(left), Inches(top), width=Inches(width_in))


# ---------------------------------------------------------------- construção
def build(cidade, sem_secao):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    gab = []  # erros injetados
    ctrl = []  # controles negativos (não devem gerar achado)

    # 1 — CAPA
    s = add_slide(prs, f"Estudo Vocacional — {cidade}")
    add_text(s, f"Análise de vocação de terreno · {cidade} · uso interno", size=20)

    # 2 — IDENTIFICACAO (raio canônico)
    s = add_slide(prs, "Identificação do Terreno e Zona de Influência")
    add_legend(s, ["Isócronas de deslocamento:", "Até 10 min", "Até 20 min", "Até 30 min"])

    # 3 — IDENTIFICACAO (raio canônico repetido -> fixa o canônico)
    s = add_slide(prs, "Zona de Influência — Deslocamento por Tempo")
    add_legend(s, ["Legenda:", "Até 10 min", "Até 20 min", "Até 30 min"])

    # 4 — IDENTIFICACAO (ERRO RADII: 45 min é estranho ao canônico {10,20,30})
    s = add_slide(prs, "Z.I. — Isócronas Consolidadas")
    add_legend(s, ["Legenda:", "Até 10 min", "Até 20 min", "Até 45 min"])
    gab.append(dict(slide=4, type="RADII", layer="DET",
                    detail="Raio '45 min' estranho ao conjunto canônico {10,20,30} min"))

    # 5 — SOCIO (ERRO ABSOLUTE_SUM nativa + ERRO SPELLING)
    s = add_slide(prs, "Perfil Sociodemográfico — Faixas de Renda")
    add_native_table(
        s,
        ["Faixa de Renda", "Domicilios"],
        [["Até 3 SM", "120"], ["3 a 6 SM", "340"], ["6 a 10 SM", "210"]],
        total_row=["Total", "700"],  # soma real = 670
        top=1.7,
    )
    add_text(s,
             "O empreedimento analisado apresenta forte demandda por unidades de dois "
             "dormitórios, com perfil de moradya compatível com a renda média local.",
             top=4.6, size=17)
    gab.append(dict(slide=5, type="ABSOLUTE_SUM", layer="DET",
                    detail="Tabela nativa: coluna 'Domicilios' soma 670 != Total 700"))
    gab.append(dict(slide=5, type="SPELLING", layer="IA",
                    detail="Erros: 'empreedimento', 'demandda', 'moradya'"))

    # 6 — SOCIO (CONTROLE: tabela-imagem CORRETA com Absoluto + % — exercita o
    # cross-check %↔absoluto e o escalonamento p/ 4o se a visão ler algum dígito errado)
    s = add_slide(prs, "Sociodemográfico — Domicílios por Faixa")
    add_table_image(
        s,
        ["Faixa", "Absoluto", "%"],
        [["Até 3 SM", "150", "21,4%"], ["3 a 6 SM", "300", "42,9%"], ["6 a 10 SM", "250", "35,7%"]],
        total_row=["Total", "700", "100%"],  # 150+300+250 = 700; %s batem com absoluto/total
    )
    ctrl.append(dict(slide=6, nota="tabela-imagem Absoluto+% correta — não deve gerar ABSOLUTE_SUM nem PERCENTAGE_SUM; se a visão errar um dígito, o cross-check %↔abs pega e escala p/ 4o"))

    # 7 — SOCIO (ERRO ABSOLUTE_SUM em imagem: coluna 2030 não fecha)
    s = add_slide(prs, "Projeção de Demanda — Renda por Ano")
    add_table_image(
        s,
        ["Faixa", "2025", "2030"],
        [["Até 3 SM", "100", "120"], ["3 a 6 SM", "200", "250"], ["6 a 10 SM", "150", "180"]],
        total_row=["Total", "450", "600"],  # 2025 ok (450); 2030 = 550 != 600
    )
    gab.append(dict(slide=7, type="ABSOLUTE_SUM", layer="VISION",
                    detail="Tabela-imagem: coluna '2030' soma 550 != Total 600"))

    # 8 — MERCADO (CONTROLE tabela nativa correta + ERRO CITY_NAME)
    s = add_slide(prs, "Mercado — Oferta Lançada")
    add_native_table(
        s,
        ["Tipologia", "Unidades"],
        [["1 dorm", "80"], ["2 dorm", "160"], ["3 dorm", "60"]],
        total_row=["Total", "300"],  # 80+160+60 = 300 OK
        top=1.7,
    )
    add_text(s,
             "A oferta lançada em Florianópolis mostra concentração no segmento médio-alto, "
             "com boa velocidade de vendas nos últimos trimestres.",
             top=4.6, size=17)
    ctrl.append(dict(slide=8, nota="tabela nativa correta — deve entrar em 'tabelas que fecham'"))
    gab.append(dict(slide=8, type="CITY_NAME", layer="IA",
                    detail=f"Menção a 'Florianópolis' em estudo de {cidade}"))

    # 9 — MERCADO (ERRO BINNING em imagem: furo entre 6000 e 9001)
    s = add_slide(prs, "Mercado — Preço Médio por Faixa (R$/m²)")
    add_table_image(
        s,
        ["Tipologia", "0-3000", "3001-6000", "9001-12000", "Acima de 12000"],
        [["Studio", "18", "10", "3", "1"], ["1 dorm", "12", "20", "8", "2"],
         ["2 dorm", "5", "18", "10", "4"], ["3 dorm", "2", "9", "12", "6"]],
    )
    gab.append(dict(slide=9, type="BINNING_RULE", layer="VISION",
                    detail="Faixas de R$/m² pulam 6001-9000 (furo entre 6000 e 9001)"))

    # 10 — LACUNAS (ERRO COHERENCE: prosa 35% vs tabela 53%)
    s = add_slide(prs, "Lacunas de Mercado")
    add_native_table(
        s,
        ["Segmento", "Participação"],
        [["Renda baixa", "22%"], ["Renda média", "25%"], ["Renda alta", "53%"]],
        top=1.7,
    )
    add_text(s,
             "Aproximadamente 35% dos domicílios concentram-se na faixa de maior renda, "
             "indicando espaço para produto de alto padrão.",
             top=4.4, size=17)
    gab.append(dict(slide=10, type="COHERENCE", layer="IA",
                    detail="Texto diz 35% mas a tabela do slide mostra 53% na renda alta"))

    # 11 — ABSORCAO (limpo)
    s = add_slide(prs, "Absorção de Mercado")
    add_text(s, "A absorção projetada acompanha o ritmo histórico da região, sem distorções.")

    # 12 — ENTORNO (ERRO LEFTOVER_NOTE)
    s = add_slide(prs, "Entorno — Mapeamento Físico e Infraestrutura")
    add_text(s, "O entorno oferece infraestrutura urbana consolidada e acesso a serviços.")
    add_text(s, "conferir dados de infraestrutura com a equipe", top=4.5, size=14)
    gab.append(dict(slide=12, type="LEFTOVER_NOTE", layer="DET",
                    detail="Nota interna 'conferir dados de infraestrutura com a equipe'"))

    # 13 — CONCLUSAO (limpo) — omitido se --sem-secao conclusao
    if sem_secao != "conclusao":
        s = add_slide(prs, "Vocação do Terreno — Conclusão e Recomendações")
        add_text(s, "O terreno apresenta vocação para produto residencial vertical de 2 dorm.")

    # 14/15 — filler limpos (sem seção canônica)
    for titulo in ("Metodologia do Estudo", "Equipe e Créditos"):
        s = add_slide(prs, titulo)
        add_text(s, "Conteúdo de apoio, sem dados numéricos auditáveis.")

    if sem_secao:
        gab.append(dict(slide=0, type="STRUCTURE_MISSING", layer="DET",
                        detail=f"Seção '{sem_secao}' ausente (--sem-secao)"))

    return prs, dict(cidade=cidade, erros=gab, controles_negativos=ctrl)


# ---------------------------------------------------------------- fixture da ata
def render_ata_png(cidade):
    """Print de página inteira da ata (>= 900x600, >= 60 KB) p/ testar o localizador."""
    W, H = 1000, 1400
    img = Image.new("RGB", (W, H), "white")
    d = ImageDraw.Draw(img)
    ft = _load_font(30)
    fb = _load_font(22)
    fh = _load_font(24)
    y = 40
    d.text((40, y), f"3. Tancredo @Kenuy / Brain", fill=(20, 40, 120), font=ft); y += 60
    blocos = [
        ("Produto pretendido:", [
            "Empreendimento com 8 torres, 1.912 unidades",
            "Todas as unidades de 2 dorm com e sem sacada.",
            "Unidades de 36m² a 41m²",
            "50% das unidades com vaga",
            "Empreendimento com características mais próximas do MCMV",
        ]),
        ("Terreno e localização:", [
            f"Av. Presidente Tancredo de Almeida Neves, bairro São Roque | {cidade}/SP",
            "Área com 25.000m² privativos",
        ]),
        ("Preço aproximado de viabilidade R$ 8.500/m²", []),
        ("Dúvida do cliente:", [
            "Qual o mix ideal de metragens.",
            "A necessidade de mais vagas",
            "Se necessário é possível fasear",
        ]),
        ("Analista da Rebrain:", [
            "Trazer tabela de lacunas de vagas de garagem",
            "Trazer análise de áreas de lazer da concorrência",
        ]),
    ]
    for titulo, itens in blocos:
        d.text((40, y), titulo, fill=(0, 0, 0), font=fh); y += 40
        for it in itens:
            d.text((70, y), f"• {it}", fill=(30, 30, 30), font=fb); y += 34
        y += 20
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def build_ata_fixture(cidade):
    """Deck mínimo (capa + ATA-imagem no slide 2 + filler) p/ testar findAtaImage."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    add_slide(prs, f"Estudo Vocacional — {cidade}")
    s = add_slide(prs, "Ata de Abertura — Briefing")
    s.shapes.add_picture(render_ata_png(cidade), Inches(0.4), Inches(0.4), height=Inches(6.8))
    s = add_slide(prs, "Metodologia")
    add_text(s, "Conteúdo de apoio.")
    return prs


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--cidade", default="Itajaí")
    ap.add_argument("--out", default="teste_corretor.pptx")
    ap.add_argument("--sem-secao", default=None,
                    help="omite uma seção canônica (ex.: entorno) p/ testar STRUCTURE_MISSING")
    ap.add_argument("--fixture", default=None, choices=["ata"],
                    help="gera um fixture específico (ata = deck mínimo com ata-imagem)")
    args = ap.parse_args()

    if args.fixture == "ata":
        prs = build_ata_fixture(args.cidade)
        prs.save(args.out)
        print(f"OK  {args.out}  ({os.path.getsize(args.out)/1e6:.2f} MB, {len(prs.slides)} slides) — fixture ata")
        return

    prs, gabarito = build(args.cidade, args.sem_secao)
    out = args.out
    prs.save(out)
    gab_path = os.path.splitext(out)[0] + ".gabarito.json"
    gabarito["arquivo"] = os.path.basename(out)
    with open(gab_path, "w", encoding="utf-8") as f:
        json.dump(gabarito, f, ensure_ascii=False, indent=2)

    size_mb = os.path.getsize(out) / 1e6
    n_det = sum(1 for e in gabarito["erros"] if e["layer"] == "DET")
    n_ia = sum(1 for e in gabarito["erros"] if e["layer"] == "IA")
    n_vis = sum(1 for e in gabarito["erros"] if e["layer"] == "VISION")
    print(f"OK  {out}  ({size_mb:.2f} MB, {len(prs.slides)} slides)")
    print(f"    gabarito: {gab_path}")
    print(f"    erros injetados: {n_det} DET (grátis) · {n_ia} IA · {n_vis} VISÃO")
    print(f"    controles negativos: {len(gabarito['controles_negativos'])}")


if __name__ == "__main__":
    main()
